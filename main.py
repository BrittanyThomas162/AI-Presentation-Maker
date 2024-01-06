import openai
from pptx.util import Pt    # used for specifying font sizes in PowerPoint files
from pptx.dml.color import RGBColor     # used to represent RGB colors for various elements in PowerPoint, such as text or shapes
from pptx import Presentation
from dotenv import load_dotenv      # used for loading environment variables from a .env file 
import os

load_dotenv()   # reads the values from a .env file and adds them to the environment variables in the current Python session

openai.api_key = os.getenv('OPENAI_API_KEY')    # sets the API key for OpenAI by retrieving its value from the environment variables

TITLE_FONT_SIZE = Pt(30)
SLIDE_FONT_SIZE = Pt(16)

def create_slide_titles(topic, num_slides):
    prompt = f"Generate {num_slides} short slide titles for the topic '{topic}'."
    response = openai.ChatCompletion.create(
        model="gpt-3.5-turbo",
        messages=[
            {"role": "system", "content": "You are a helpful assistant."},
            {"role": "user", "content": prompt}
        ],
        temperature=0.0,
        top_p=0.1,
        max_tokens=200,
        n=1
    )
    return response['choices'][0]['message']['content'].split("\n")

def create_slide_content(slide_title):
    prompt = f"Generate content for the slide: '{slide_title}'. The content must be in medium-worded paragraphs. Only return 2 paragraphs."
    response = openai.ChatCompletion.create(
        model="gpt-3.5-turbo",
        messages=[
            {"role": "system", "content": "You are a helpful assistant."},
            {"role": "user", "content": prompt}
        ],
        temperature=0.0,
        top_p=0.1,
        max_tokens=300,
        n=1
    )
    return response['choices'][0]['message']['content']

def create_presentation(topic, slide_titles, slide_contents):
    powerpoint = Presentation()     # used to create title and content slide, change font colors and edits
    
    title_slide_layout = powerpoint.slide_layouts[0]
    content_slide_layout = powerpoint.slide_layouts[1]

    background_color = RGBColor(173, 216, 230)

    title_slide = powerpoint.slides.add_slide(title_slide_layout)
    title = title_slide.shapes.title
    title.text = topic

    title.text_frame.paragraphs[0].font.size = Pt(48)
    title.text_frame.paragraphs[0].font.bold = True
    content = title_slide.placeholders[1]       # 0 index is for the titles and 1 index is for the content
    content.text = "Created by AI."
    content.text_frame.paragraphs[0].font.size = Pt(24)
    
    background = title_slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = background_color

    for slide_title, slide_content in zip(slide_titles, slide_contents):
        
        slide = powerpoint.slides.add_slide(content_slide_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = background_color

        title = slide.shapes.title
        title.text = slide_title
        title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
        title.text_frame.paragraphs[0].font.bold = True

        content = slide.placeholders[1]
        content.text = slide_content 
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = SLIDE_FONT_SIZE
    

    powerpoint.save(f"powerpoints/{topic}.pptx")


def main():
    topic = "AI in Autonomous Vehicles"
    num_slides = 5

    slide_titles = create_slide_titles(topic, num_slides)
    print("Generated Slide Titles.")
   
    filtered_slide_titles =  [item for item in slide_titles if item.strip() != '']
    slide_contents = [create_slide_content(title) for title in filtered_slide_titles]
    print("Generated Slide Content.")
    create_presentation(topic, filtered_slide_titles, slide_contents)
    
    print("The presentation is generated!")


if __name__ == "__main__":
    main()







'''from openai import OpenAI
from pptx.util import Pt    # used for specifying font sizes in PowerPoint files
from pptx.dml.color import RGBColor     # used to represent RGB colors for various elements in PowerPoint, such as text or shapes
from pptx import Presentation
#from dotenv import load_dotenv      # used for loading environment variables from a .env file 
import os

client = OpenAI()
#load_dotenv()   # reads the values from a .env file and adds them to the environment variables in the current Python session
#client.api_key = os.getenv('OPENAI_API_KEY')    # sets the API key for OpenAI by retrieving its value from the environment variables

TITLE_FONT_SIZE = Pt(30)
SLIDE_FONT_SIZE = Pt(16)


def create_slide_titles(topic, num_slides):
    prompt = f"Generate {num_slides} short slide titles for the topic '{topic}'."
    response = client.chat.completions.create (
        model="gpt-3.5-turbo",
        messages=[{"role":"system","content":prompt}],
        temperature=0.0,
        top_p=0.1,
        max_tokens=200,
        n=1,
        timeout = 15)
    return response.choices[0].messages.content.split("\n")

def create_slide_content(slide_title):
    prompt = f"Generate content for the slide: '{slide_title}'. The content must be in medium worded paragraphs. Only return 2 paragraphs."
    response = client.chat.completions.create  (
        model="gpt-3.5-turbo",
        messages=[{"role":"system","content":prompt}],
        temperature=0.0,
        top_p=0.1,
        max_tokens=300,
        n=1,
        timeout = 15
    )
    return response.choices[0].messages.content

def create_presentation(topic, slide_titles, slide_contents):
    powerpoint = Presentation()     # used to create title and content slide, change font colors and edits
    
    title_slide_layout = powerpoint.slide_layouts[0]
    content_slide_layout = powerpoint.slide_layouts[1]

    background_color = RGBColor(173, 216, 230)

    title_slide = powerpoint.slides.add_slide(title_slide_layout)
    title = title_slide.shapes.title
    title.text = topic

    title.text_frame.paragraphs[0].font.size = Pt(48)
    title.text_frame.paragraphs[0].font.bold = True
    content = title_slide.placeholders[1]       # 0 index is for the titles and 1 index is for the content
    content.text = "Created by AI."
    content.text_frame.paragraphs[0].font.size = Pt(24)
    
    background = title_slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = background_color

    for slide_title, slide_content in zip(slide_titles, slide_contents):
        
        slide = powerpoint.slides.add_slide(content_slide_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = background_color

        title = slide.shapes.title
        title.text = slide_title
        title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
        title.text_frame.paragraphs[0].font.bold = True

        content = slide.placeholders[1]
        content.text = slide_content 
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = SLIDE_FONT_SIZE
    

    powerpoint.save(f"powerpoints/{topic}.pptx")


def main():
    topic = "AI in Autonomous Vehicles"
    num_slides = 5

    slide_titles = create_slide_titles(topic, num_slides)
    print("Generated Slide Titles.")
   
    filtered_slide_titles =  [item for item in slide_titles if item.strip() != '']
    slide_contents = [create_slide_content(title) for title in filtered_slide_titles]
    print("Generated Slide Content.")
    create_presentation(topic, filtered_slide_titles, slide_contents)
    
    print("The Presentation is generated!")


if __name__ == "__main__":
    main()

'''